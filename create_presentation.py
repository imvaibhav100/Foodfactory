#!/usr/bin/env python3
"""
Script to create a professional PowerPoint presentation for the FoodFactory project.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_foodfactory_presentation(output_dir=None):
    """Create a professional presentation for FoodFactory project."""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (Modern food-themed colors)
    PRIMARY_COLOR = RGBColor(230, 74, 25)  # Orange-red
    SECONDARY_COLOR = RGBColor(255, 165, 0)  # Orange
    ACCENT_COLOR = RGBColor(34, 139, 34)  # Green
    TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray
    LIGHT_BG = RGBColor(255, 248, 240)  # Light cream
    
    # =====================================================
    # SLIDE 1: TITLE SLIDE
    # =====================================================
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "FoodFactory 🍽️"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Modern Full-Stack Food Ordering Web Application"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = TEXT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Author info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    author_frame = author_box.text_frame
    author_frame.text = "by Vaibhav"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(20)
    author_para.font.color.rgb = SECONDARY_COLOR
    author_para.alignment = PP_ALIGN.CENTER
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(3), Inches(5.5), Inches(4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_COLOR
    line.line.fill.background()
    
    # =====================================================
    # SLIDE 2: INTRODUCTION
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Introduction"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line under header
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Content paragraphs
    intro_text = [
        "FoodFactory is a modern, full-stack food ordering web application designed to revolutionize the online food delivery experience.",
        "",
        "Built with enterprise-grade technologies, FoodFactory provides:",
        "• A seamless platform for users to browse and order delicious food online",
        "• Robust admin management capabilities for efficient operations",
        "• Secure authentication and user management",
        "• Beautiful, responsive design for optimal user experience",
        "",
        "This project demonstrates real-world application development with industry best practices."
    ]
    
    for i, text in enumerate(intro_text):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
        if text and not text.startswith("•"):
            p.font.bold = True if i == 2 else False
    
    # =====================================================
    # SLIDE 3: KEY FEATURES
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Key Features"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Left column features
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(4), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    left_features = [
        "👤 User Management",
        "• Secure user registration & login",
        "• Password visibility toggle",
        "• Role-based access control",
        "",
        "🍕 Product Browsing",
        "• Wide variety of food items",
        "• Intuitive product catalog",
        "• Category-based organization",
        "",
        "🛒 Order Management",
        "• Seamless order placement",
        "• Order success confirmation",
        "• User-friendly checkout process"
    ]
    
    for i, text in enumerate(left_features):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•") and text != ""
        p.space_after = Pt(6)
    
    # Right column features
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    right_features = [
        "⚙️ Admin Dashboard",
        "• Complete user management",
        "• Product CRUD operations",
        "• Admin account management",
        "",
        "🎨 Modern UI/UX",
        "• Responsive design",
        "• Beautiful CSS styling",
        "• Mobile-friendly interface",
        "",
        "🔒 Security Features",
        "• Spring Security integration",
        "• Secure authentication",
        "• Data validation & error handling"
    ]
    
    for i, text in enumerate(right_features):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•") and text != ""
        p.space_after = Pt(6)
    
    # =====================================================
    # SLIDE 4: TECHNOLOGY STACK
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Technology Stack"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Backend section
    backend_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(4), Inches(2))
    backend_box.fill.solid()
    backend_box.fill.fore_color.rgb = LIGHT_BG
    backend_frame = backend_box.text_frame
    backend_frame.word_wrap = True
    
    backend_text = [
        "💻 Backend Technologies",
        "• Java 17",
        "• Spring Boot 3.1.3",
        "• Spring Security",
        "• Spring Data JPA",
        "• Maven"
    ]
    
    for i, text in enumerate(backend_text):
        if i == 0:
            p = backend_frame.paragraphs[0]
        else:
            p = backend_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•")
        p.space_after = Pt(6)
    
    # Frontend section
    frontend_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.2), Inches(4), Inches(2))
    frontend_box.fill.solid()
    frontend_box.fill.fore_color.rgb = LIGHT_BG
    frontend_frame = frontend_box.text_frame
    frontend_frame.word_wrap = True
    
    frontend_text = [
        "🎨 Frontend Technologies",
        "• Thymeleaf Templates",
        "• HTML5",
        "• CSS3",
        "• JavaScript",
        "• Responsive Design"
    ]
    
    for i, text in enumerate(frontend_text):
        if i == 0:
            p = frontend_frame.paragraphs[0]
        else:
            p = frontend_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•")
        p.space_after = Pt(6)
    
    # Database section
    database_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(4), Inches(1.5))
    database_box.fill.solid()
    database_box.fill.fore_color.rgb = LIGHT_BG
    database_frame = database_box.text_frame
    database_frame.word_wrap = True
    
    database_text = [
        "🗄️ Database",
        "• MySQL",
        "• JPA/Hibernate ORM",
        "• JDBC Connection"
    ]
    
    for i, text in enumerate(database_text):
        if i == 0:
            p = database_frame.paragraphs[0]
        else:
            p = database_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•")
        p.space_after = Pt(6)
    
    # Tools section
    tools_box = slide.shapes.add_textbox(Inches(5.5), Inches(4.5), Inches(4), Inches(1.5))
    tools_box.fill.solid()
    tools_box.fill.fore_color.rgb = LIGHT_BG
    tools_frame = tools_box.text_frame
    tools_frame.word_wrap = True
    
    tools_text = [
        "🛠️ Development Tools",
        "• Git & GitHub",
        "• Spring Boot DevTools",
        "• Maven Build System"
    ]
    
    for i, text in enumerate(tools_text):
        if i == 0:
            p = tools_frame.paragraphs[0]
        else:
            p = tools_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•")
        p.space_after = Pt(6)
    
    # =====================================================
    # SLIDE 5: ARCHITECTURE & DESIGN
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Architecture & Design Patterns"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(3.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    arch_text = [
        "🏗️ MVC Architecture Pattern",
        "• Model: Entity classes with JPA annotations",
        "• View: Thymeleaf templates for dynamic rendering",
        "• Controller: RESTful controllers handling HTTP requests",
        "",
        "📦 Layered Architecture",
        "• Controller Layer: Request handling and routing",
        "• Service Layer: Business logic implementation",
        "• Repository Layer: Data access with Spring Data JPA",
        "• Entity Layer: Domain models and database mapping",
        "",
        "🎯 Key Design Principles",
        "• Separation of Concerns",
        "• Dependency Injection",
        "• Single Responsibility Principle",
        "• RESTful API design conventions"
    ]
    
    for i, text in enumerate(arch_text):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•") and text != ""
        p.space_after = Pt(8)
    
    # =====================================================
    # SLIDE 6: CHALLENGES FACED
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Challenges Faced & Solutions"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(3.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Challenge boxes
    challenges = [
        {
            "title": "🔐 Security Implementation",
            "challenge": "Challenge: Implementing secure authentication and authorization",
            "solution": "Solution: Integrated Spring Security with custom user details service and password encoding"
        },
        {
            "title": "🎨 UI/UX Consistency",
            "challenge": "Challenge: Maintaining consistent design across multiple pages",
            "solution": "Solution: Created modular CSS files and reusable Thymeleaf templates"
        },
        {
            "title": "💾 Database Integration",
            "challenge": "Challenge: Managing complex entity relationships",
            "solution": "Solution: Utilized JPA annotations and Spring Data repositories for efficient data handling"
        }
    ]
    
    y_pos = 2.2
    for challenge_data in challenges:
        # Challenge box
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        frame = box.text_frame
        frame.word_wrap = True
        
        # Title
        p = frame.paragraphs[0]
        p.text = challenge_data["title"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.space_after = Pt(6)
        
        # Challenge
        p = frame.add_paragraph()
        p.text = challenge_data["challenge"]
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(4)
        
        # Solution
        p = frame.add_paragraph()
        p.text = challenge_data["solution"]
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_COLOR
        p.font.bold = True
        
        y_pos += 1.5
    
    # =====================================================
    # SLIDE 7: PROJECT HIGHLIGHTS
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Project Highlights"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2.8), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Highlights in boxes
    highlights = [
        ("✨", "Production-Ready", "Enterprise-grade code with proper error handling and validation"),
        ("🚀", "Scalable Architecture", "Modular design that supports future enhancements and features"),
        ("📱", "Responsive Design", "Works seamlessly across desktop, tablet, and mobile devices"),
        ("🔄", "Real-Time Operations", "Dynamic content updates and seamless user interactions")
    ]
    
    x_positions = [Inches(0.8), Inches(5.1), Inches(0.8), Inches(5.1)]
    y_positions = [Inches(2.2), Inches(2.2), Inches(4.2), Inches(4.2)]
    
    for i, (emoji, title, desc) in enumerate(highlights):
        box = slide.shapes.add_textbox(x_positions[i], y_positions[i], Inches(3.9), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        frame = box.text_frame
        frame.word_wrap = True
        
        # Emoji
        p = frame.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)
        
        # Title
        p = frame.add_paragraph()
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)
        
        # Description
        p = frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # =====================================================
    # SLIDE 8: FUTURE ENHANCEMENTS
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Future Enhancements"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2.8), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    future_text = [
        "🔮 Planned Features & Improvements",
        "",
        "💳 Payment Gateway Integration",
        "• Support for multiple payment methods",
        "• Secure online transactions",
        "",
        "📍 Real-Time Order Tracking",
        "• Live order status updates",
        "• GPS-based delivery tracking",
        "",
        "🤖 AI-Powered Recommendations",
        "• Personalized food suggestions",
        "• Smart ordering based on preferences",
        "",
        "📊 Analytics Dashboard",
        "• Sales reports and insights",
        "• User behavior analytics",
        "",
        "🌐 Multi-Language Support",
        "• Internationalization capabilities"
    ]
    
    for i, text in enumerate(future_text):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("•") else Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = not text.startswith("•") and text != ""
        p.space_after = Pt(6)
    
    # =====================================================
    # SLIDE 9: CONCLUSION
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header.text_frame
    header_frame.text = "Conclusion"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(44)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Decorative line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()
    
    # Content box with gradient-like effect
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(3.5))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = LIGHT_BG
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    conclusion_text = [
        "FoodFactory successfully demonstrates the implementation of a full-featured, production-ready food ordering platform using modern web technologies.",
        "",
        "Key Takeaways:",
        "✓ Robust backend with Spring Boot and Spring Security",
        "✓ Intuitive frontend with Thymeleaf and responsive design",
        "✓ Scalable architecture following industry best practices",
        "✓ Complete CRUD operations for users, admins, and products",
        "",
        "This project showcases practical application of full-stack development skills and serves as a solid foundation for real-world food delivery services."
    ]
    
    for i, text in enumerate(conclusion_text):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.startswith("✓") else Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = text.startswith("Key Takeaways")
        p.space_after = Pt(8)
        if text.startswith("✓"):
            p.font.color.rgb = ACCENT_COLOR
    
    # Call to action
    cta_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.6))
    cta_frame = cta_box.text_frame
    cta_frame.text = "🔗 GitHub: github.com/imvaibhav100/Foodfactory"
    cta_para = cta_frame.paragraphs[0]
    cta_para.font.size = Pt(20)
    cta_para.font.color.rgb = PRIMARY_COLOR
    cta_para.font.bold = True
    cta_para.alignment = PP_ALIGN.CENTER
    
    # =====================================================
    # SLIDE 10: Q&A
    # =====================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Large Q&A text
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    qa_frame = qa_box.text_frame
    qa_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    qa_frame.text = "Questions & Answers"
    qa_para = qa_frame.paragraphs[0]
    qa_para.font.size = Pt(60)
    qa_para.font.bold = True
    qa_para.font.color.rgb = PRIMARY_COLOR
    qa_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Thank you for your attention! 🍽️"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = TEXT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Feel free to ask any questions!"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(18)
    contact_para.font.color.rgb = SECONDARY_COLOR
    contact_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    if output_dir is None:
        # Use the directory where the script is located
        script_dir = os.path.dirname(os.path.abspath(__file__))
        output_path = os.path.join(script_dir, 'FoodFactory_Presentation.pptx')
    else:
        output_path = os.path.join(output_dir, 'FoodFactory_Presentation.pptx')
    
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    create_foodfactory_presentation()
